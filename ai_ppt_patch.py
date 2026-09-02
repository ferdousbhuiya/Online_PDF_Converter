"""AI-first PDF to PowerPoint conversion for PDFMaster Pro.

AI Presentation mode extracts PDF text locally, asks Groq for a structured slide
plan, then creates an editable .pptx with python-pptx. Preserve Layout mode
continues to use the existing fidelity converter as a fallback/alternative.
"""

import json
import os
import re
import urllib.error
import urllib.request

from converter_patch import handle_pdf_to_ppt as handle_preserve_layout_ppt


GROQ_CHAT_URL = "https://api.groq.com/openai/v1/chat/completions"
DEFAULT_GROQ_MODEL = os.environ.get("GROQ_MODEL", "llama-3.3-70b-versatile")


def _clean_text(value):
    return re.sub(r"\s+", " ", (value or "")).strip()


def _extract_pdf_content(path):
    import pdfplumber

    pages = []
    total_chars = 0
    with pdfplumber.open(path) as pdf:
        for index, page in enumerate(pdf.pages, start=1):
            text = page.extract_text(x_tolerance=2, y_tolerance=3) or ""
            text = text.strip()
            tables = []
            try:
                for table in page.extract_tables()[:3]:
                    cleaned_rows = []
                    for row in table or []:
                        cleaned = [_clean_text(cell) for cell in (row or [])]
                        if any(cleaned):
                            cleaned_rows.append(cleaned)
                    if cleaned_rows:
                        tables.append(cleaned_rows[:20])
            except Exception:
                pass

            pages.append({
                "page": index,
                "text": text,
                "tables": tables,
            })
            total_chars += len(text)

    if not pages:
        raise ValueError("The PDF contains no pages")
    if total_chars < 40:
        raise ValueError("This PDF has too little extractable text for AI presentation mode. Use Preserve Layout for scanned/image-only PDFs, or run OCR first.")
    return pages


def _target_slide_count(page_count, length_mode):
    length_mode = (length_mode or "standard").lower()
    if length_mode == "short":
        return max(4, min(10, round(page_count * 0.45) + 2))
    if length_mode == "detailed":
        return max(7, min(28, round(page_count * 0.95) + 2))
    return max(5, min(18, round(page_count * 0.65) + 2))


def _build_source_packet(pages, max_chars=52000):
    chunks = []
    used = 0
    for page in pages:
        text = page["text"]
        table_text = ""
        if page["tables"]:
            table_lines = []
            for t_index, table in enumerate(page["tables"], start=1):
                table_lines.append(f"Table {t_index}:")
                table_lines.extend(" | ".join(row) for row in table[:12])
            table_text = "\n" + "\n".join(table_lines)
        block = f"\n--- SOURCE PAGE {page['page']} ---\n{text}{table_text}\n"
        if used + len(block) > max_chars:
            remaining = max_chars - used
            if remaining > 600:
                chunks.append(block[:remaining])
            break
        chunks.append(block)
        used += len(block)
    return "".join(chunks)


def _groq_slide_plan(pages, length_mode, style_mode, include_tables):
    api_key = os.environ.get("GROQ_API_KEY", "").strip()
    if not api_key:
        raise RuntimeError("AI Presentation requires GROQ_API_KEY in the server environment. Add the same Groq API key pattern used by Skills Pathfinder to Coolify, then redeploy.")

    page_count = len(pages)
    target = _target_slide_count(page_count, length_mode)
    source = _build_source_packet(pages)
    style_mode = (style_mode or "professional").lower()

    system_prompt = (
        "You are an expert presentation designer. Convert source-document content into a real presentation, not page screenshots. "
        "Create a coherent narrative, concise editable slide text, meaningful slide divisions, and accurate source-page references. "
        "Do not invent facts. Return strict JSON only."
    )

    user_prompt = f"""
Create an editable PowerPoint slide plan from the PDF content below.

Presentation length: {length_mode}
Target slide count: approximately {target}
Presentation style: {style_mode}
Include useful source tables: {str(bool(include_tables)).lower()}
Source PDF pages: {page_count}

Requirements:
- Build a proper presentation, not one slide per PDF page.
- Start with a strong title slide unless the content is clearly unsuitable.
- Group related ideas into logical sections.
- Use short presentation-ready bullets, normally 3-6 bullets per content slide.
- Split dense source material across multiple slides when useful.
- Combine short related material when useful.
- Use source page numbers for traceability.
- When a source table contains useful information, create a table slide or a slide with a small table.
- Do not copy long paragraphs verbatim.
- Do not invent numerical values, conclusions, names, or citations not present in the source.
- End with a conclusion/summary slide when appropriate.

Return exactly this JSON structure:
{{
  "deck_title": "...",
  "deck_subtitle": "...",
  "slides": [
    {{
      "title": "...",
      "layout": "title|bullets|two_column|table|section|conclusion",
      "bullets": ["..."],
      "left_bullets": ["..."],
      "right_bullets": ["..."],
      "table": {{"headers": ["..."], "rows": [["..."]]}},
      "source_pages": [1,2],
      "speaker_notes": "optional concise context"
    }}
  ]
}}

SOURCE CONTENT:
{source}
"""

    body = {
        "model": DEFAULT_GROQ_MODEL,
        "temperature": 0.2,
        "response_format": {"type": "json_object"},
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
    }
    request = urllib.request.Request(
        GROQ_CHAT_URL,
        data=json.dumps(body).encode("utf-8"),
        headers={
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json",
            "User-Agent": "PDFMaster-Pro/1.0",
        },
        method="POST",
    )
    try:
        with urllib.request.urlopen(request, timeout=120) as response:
            payload = json.loads(response.read().decode("utf-8"))
    except urllib.error.HTTPError as exc:
        try:
            detail = exc.read().decode("utf-8", errors="ignore")[:600]
        except Exception:
            detail = ""
        raise RuntimeError(f"Groq API request failed ({exc.code}). {detail}") from exc
    except urllib.error.URLError as exc:
        raise RuntimeError(f"Could not reach Groq API: {exc.reason}") from exc

    content = payload.get("choices", [{}])[0].get("message", {}).get("content", "")
    if not content:
        raise RuntimeError("Groq returned an empty slide plan")
    try:
        plan = json.loads(content)
    except json.JSONDecodeError:
        match = re.search(r"\{.*\}", content, re.S)
        if not match:
            raise RuntimeError("Groq did not return valid presentation JSON")
        plan = json.loads(match.group(0))

    slides = plan.get("slides")
    if not isinstance(slides, list) or not slides:
        raise RuntimeError("Groq returned a presentation plan without slides")
    return plan


def _render_source_pages(src, output_dir, requested_pages):
    """Render only requested PDF pages as optional visual references."""
    from pdf2image import convert_from_path

    rendered = {}
    for page_number in sorted(set(p for p in requested_pages if isinstance(p, int) and p > 0)):
        try:
            images = convert_from_path(src, dpi=115, fmt="jpeg", first_page=page_number, last_page=page_number, thread_count=1)
            if not images:
                continue
            path = os.path.join(output_dir, f"ai_source_page_{page_number}.jpg")
            images[0].save(path, "JPEG", quality=82, optimize=True)
            rendered[page_number] = path
        except Exception:
            continue
    return rendered


def _add_text_box(slide, text, left, top, width, height, font_size=22, bold=False, color=None, align=None):
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
    from pptx.util import Pt

    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    if align:
        paragraph.alignment = align
    run = paragraph.add_run()
    run.text = _clean_text(text)
    run.font.name = "Aptos"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*(color or (24, 32, 48)))
    return box


def _add_bullet_box(slide, bullets, left, top, width, height, font_size=22):
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    bullets = [_clean_text(x) for x in (bullets or []) if _clean_text(x)][:7]
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    if not bullets:
        bullets = ["Content summarized from the source document."]
    for idx, bullet in enumerate(bullets):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(font_size)
        p.font.color.rgb = RGBColor(45, 55, 72)
        p.space_after = Pt(9)
        try:
            p._p.get_or_add_pPr().insert(0, p._p._new_buChar())
        except Exception:
            pass
    return box


def _add_footer(slide, source_pages, slide_width, slide_height):
    pages = [str(p) for p in source_pages or [] if isinstance(p, int)]
    if not pages:
        return
    from pptx.util import Inches
    _add_text_box(
        slide,
        "Source PDF pages: " + ", ".join(pages[:8]),
        Inches(0.55), slide_height - Inches(0.34), slide_width - Inches(1.1), Inches(0.18),
        font_size=8,
        color=(110, 118, 132),
    )


def _add_table(slide, table_data, left, top, width, height):
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    headers = table_data.get("headers") if isinstance(table_data, dict) else None
    rows = table_data.get("rows") if isinstance(table_data, dict) else None
    headers = [str(x) for x in (headers or [])][:6]
    rows = [[str(x) for x in row[:6]] for row in (rows or [])[:8] if isinstance(row, list)]
    if not headers and rows:
        headers = [f"Column {i+1}" for i in range(len(rows[0]))]
    if not headers:
        return False
    cols = len(headers)
    normalized_rows = [(row + [""] * cols)[:cols] for row in rows]
    shape = slide.shapes.add_table(len(normalized_rows) + 1, cols, left, top, width, height)
    table = shape.table
    for c, value in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = value
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(255, 255, 255)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(55, 65, 81)
    for r, row in enumerate(normalized_rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(11)
                p.font.color.rgb = RGBColor(42, 48, 60)
    return True


def _build_pptx(plan, src, output_dir, style_mode, include_source_images):
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    sw, sh = prs.slide_width, prs.slide_height

    slides = plan.get("slides") or []
    requested_pages = []
    if include_source_images:
        for spec in slides:
            for page in spec.get("source_pages") or []:
                if isinstance(page, int):
                    requested_pages.append(page)
    rendered = _render_source_pages(src, output_dir, requested_pages)

    for index, spec in enumerate(slides, start=1):
        slide = prs.slides.add_slide(blank)
        layout = (spec.get("layout") or "bullets").lower()
        title = _clean_text(spec.get("title") or f"Slide {index}")

        # Background
        bg = slide.background.fill
        bg.solid()
        if style_mode == "academic":
            bg.fore_color.rgb = RGBColor(250, 250, 248)
        elif style_mode == "simple":
            bg.fore_color.rgb = RGBColor(255, 255, 255)
        else:
            bg.fore_color.rgb = RGBColor(248, 250, 252)

        # Accent bar
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.12), sh)
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(79, 70, 229)
        accent.line.fill.background()

        if layout == "title" or (index == 1 and len(slides) > 1):
            _add_text_box(slide, title, Inches(0.9), Inches(1.75), Inches(11.5), Inches(1.4), 34, True, (17, 24, 39), PP_ALIGN.CENTER)
            subtitle = _clean_text(plan.get("deck_subtitle") or spec.get("speaker_notes") or "AI-generated presentation from the source PDF")
            _add_text_box(slide, subtitle, Inches(1.5), Inches(3.3), Inches(10.3), Inches(0.9), 18, False, (90, 98, 112), PP_ALIGN.CENTER)
            _add_footer(slide, spec.get("source_pages"), sw, sh)
            continue

        _add_text_box(slide, title, Inches(0.65), Inches(0.45), Inches(12.0), Inches(0.65), 26, True, (17, 24, 39))

        source_pages = spec.get("source_pages") or []
        preview_path = None
        if include_source_images:
            preview_path = next((rendered.get(p) for p in source_pages if rendered.get(p)), None)

        if layout == "two_column":
            _add_bullet_box(slide, spec.get("left_bullets") or spec.get("bullets"), Inches(0.75), Inches(1.4), Inches(5.65), Inches(5.2), 19)
            _add_bullet_box(slide, spec.get("right_bullets"), Inches(6.8), Inches(1.4), Inches(5.65), Inches(5.2), 19)
        elif layout == "table" and _add_table(slide, spec.get("table") or {}, Inches(0.8), Inches(1.5), Inches(11.8), Inches(4.8)):
            bullets = spec.get("bullets") or []
            if bullets:
                _add_text_box(slide, "Key takeaway: " + bullets[0], Inches(0.85), Inches(6.35), Inches(11.5), Inches(0.45), 13, False, (75, 85, 99))
        elif layout in ("section", "conclusion"):
            _add_text_box(slide, title, Inches(1.0), Inches(2.0), Inches(11.2), Inches(1.0), 32, True, (17, 24, 39), PP_ALIGN.CENTER)
            _add_bullet_box(slide, spec.get("bullets"), Inches(2.0), Inches(3.2), Inches(9.3), Inches(2.6), 20)
        elif preview_path:
            _add_bullet_box(slide, spec.get("bullets"), Inches(0.75), Inches(1.45), Inches(7.25), Inches(4.95), 20)
            try:
                slide.shapes.add_picture(preview_path, Inches(8.35), Inches(1.55), width=Inches(4.25), height=Inches(4.8))
            except Exception:
                pass
        else:
            _add_bullet_box(slide, spec.get("bullets"), Inches(0.9), Inches(1.45), Inches(11.5), Inches(5.2), 21)

        _add_footer(slide, source_pages, sw, sh)

        notes = _clean_text(spec.get("speaker_notes"))
        if notes:
            try:
                notes_slide = slide.notes_slide
                notes_tf = notes_slide.notes_text_frame
                notes_tf.text = notes
            except Exception:
                pass

    if not prs.slides:
        raise RuntimeError("The AI slide plan produced no slides")

    output = os.path.join(output_dir, "converted.pptx")
    prs.save(output)
    if not os.path.exists(output) or os.path.getsize(output) < 1000:
        raise RuntimeError("PowerPoint file was not created correctly")

    for path in rendered.values():
        try:
            os.remove(path)
        except OSError:
            pass
    return output, len(prs.slides)


def handle_pdf_to_ppt(files, output_dir, form_data):
    """AI Presentation by default; Preserve Layout remains available."""
    mode = (form_data.get("ppt_mode") or "ai").lower()
    if mode == "preserve":
        return handle_preserve_layout_ppt(files, output_dir, form_data)

    try:
        if not files:
            return {"success": False, "error": "No PDF uploaded"}
        src = files[0]
        length_mode = (form_data.get("ppt_length") or "standard").lower()
        style_mode = (form_data.get("ppt_style") or "professional").lower()
        include_source_images = str(form_data.get("ppt_include_images", "yes")).lower() not in ("no", "false", "0")
        include_tables = str(form_data.get("ppt_include_tables", "yes")).lower() not in ("no", "false", "0")

        if length_mode not in ("short", "standard", "detailed"):
            length_mode = "standard"
        if style_mode not in ("professional", "academic", "simple"):
            style_mode = "professional"

        pages = _extract_pdf_content(src)
        plan = _groq_slide_plan(pages, length_mode, style_mode, include_tables)
        output, slide_count = _build_pptx(plan, src, output_dir, style_mode, include_source_images)
        return {
            "success": True,
            "download_url": f"/download/{os.path.basename(output_dir)}/{os.path.basename(output)}",
            "filename": "converted.pptx",
            "message": f"AI Presentation created {slide_count} editable slide(s) from {len(pages)} PDF page(s) using Groq. Titles, key points and slide divisions were generated from the source content.",
        }
    except Exception as exc:
        return {"success": False, "error": f"AI PDF to PowerPoint failed: {exc}"}
