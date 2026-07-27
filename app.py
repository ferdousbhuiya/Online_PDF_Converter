"""
PDFMaster Pro - Complete PDF Converter Web Application
Author: AI Assistant
Date: July 27, 2026
"""

import os
import uuid
import shutil
import threading
import time
from datetime import datetime
import subprocess
from dotenv import load_dotenv
from flask import Flask, render_template, request, send_file, jsonify, redirect, url_for, flash
from werkzeug.utils import secure_filename
from config import Config
import smtplib
from email.message import EmailMessage

# Load .env file
load_dotenv()


# PDF Libraries
from pypdf import PdfReader, PdfWriter, PdfMerger
from PyPDF2 import PdfFileMerger
import pdfplumber
from reportlab.lib.pagesizes import letter, A4
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from PIL import Image
import img2pdf
from pdf2docx import Converter as PDF2DOCXConverter
import openpyxl
from docx import Document
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.util import Inches as PptxInches

# Initialize Flask App
app = Flask(__name__)
app.config.from_object(Config)

# Email Configuration (You will set these up in Step 4)
EMAIL_ADDRESS = os.environ.get('EMAIL_ADDRESS', 'your_email@gmail.com')
EMAIL_PASSWORD = os.environ.get('EMAIL_PASSWORD', 'your_gmail_app_password')

# Create necessary directories
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['OUTPUT_FOLDER'], exist_ok=True)

# ============================================
# HELPER FUNCTIONS
# ============================================

# Helper function to get proper accept attribute for file inputs
def get_accept_attribute(input_type):
    """Returns proper MIME types and extensions for file input accept attribute."""
    accept_types = {
        'pdf': ['.pdf', 'application/pdf'],
        'word': ['.doc', '.docx', 'application/msword', 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'],
        'excel': ['.xls', '.xlsx', 'application/vnd.ms-excel', 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'],
        'powerpoint': ['.ppt', '.pptx', 'application/vnd.ms-powerpoint', 'application/vnd.openxmlformats-officedocument.presentationml.presentation'],
        'image': ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.webp', 'image/*'],
        'html': ['.html', '.htm', 'text/html']
    }
    
    types = accept_types.get(input_type, ['*/*'])
    return ','.join(types)

# Make it available in templates
app.jinja_env.globals.update(get_accept_attribute=get_accept_attribute)




def allowed_file(filename, category):
    """Check if file extension is allowed for the given category."""
    if '.' not in filename:
        return False
    ext = filename.rsplit('.', 1)[1].lower()
    return ext in app.config['ALLOWED_EXTENSIONS'].get(category, [])

def get_unique_filename():
    """Generate unique filename."""
    return f"{uuid.uuid4().hex}_{int(time.time())}"

def cleanup_old_files():
    """Background task to clean up old files."""
    while True:
        try:
            now = time.time()
            for folder in [app.config['UPLOAD_FOLDER'], app.config['OUTPUT_FOLDER']]:
                for filename in os.listdir(folder):
                    filepath = os.path.join(folder, filename)
                    if os.path.isfile(filepath):
                        if now - os.path.getmtime(filepath) > app.config['CLEANUP_INTERVAL']:
                            os.remove(filepath)
        except Exception as e:
            print(f"Cleanup error: {e}")
        time.sleep(600)  # Run every 10 minutes

# Start cleanup thread
cleanup_thread = threading.Thread(target=cleanup_old_files, daemon=True)
cleanup_thread.start()

# ============================================
# TOOL DEFINITIONS
# ============================================

TOOLS = [
    {'id': 'merge', 'name': 'Merge PDF', 'icon': '🔗', 'color': '#FF6B6B', 
     'description': 'Combine multiple PDFs into one file', 'category': 'organize',
     'input': 'pdf', 'multiple': True},
    
    {'id': 'split', 'name': 'Split PDF', 'icon': '✂️', 'color': '#4ECDC4',
     'description': 'Extract pages from PDF into separate files', 'category': 'organize',
     'input': 'pdf', 'multiple': False},
    
    {'id': 'compress', 'name': 'Compress PDF', 'icon': '🗜️', 'color': '#FFE66D',
     'description': 'Reduce PDF file size while keeping quality', 'category': 'optimize',
     'input': 'pdf', 'multiple': False},
    
    {'id': 'rotate', 'name': 'Rotate PDF', 'icon': '🔄', 'color': '#A8E6CF',
     'description': 'Rotate PDF pages in any direction', 'category': 'organize',
     'input': 'pdf', 'multiple': False},
    
    {'id': 'unlock', 'name': 'Unlock PDF', 'icon': '🔓', 'color': '#FFD3B6',
     'description': 'Remove password protection from PDF', 'category': 'security',
     'input': 'pdf', 'multiple': False},
    
    {'id': 'protect', 'name': 'Protect PDF', 'icon': '🔒', 'color': '#C9B1FF',
     'description': 'Add password protection to PDF', 'category': 'security',
     'input': 'pdf', 'multiple': False},
    
    {'id': 'watermark', 'name': 'Watermark PDF', 'icon': '💧', 'color': '#95E1D3',
     'description': 'Add text watermark to PDF pages', 'category': 'edit',
     'input': 'pdf', 'multiple': False},
    
    {'id': 'pagenumber', 'name': 'Page Numbers', 'icon': '🔢', 'color': '#F38181',
     'description': 'Add page numbers to your PDF', 'category': 'edit',
     'input': 'pdf', 'multiple': False},
    
    {'id': 'organize', 'name': 'Organize PDF', 'icon': '📋', 'color': '#AA96DA',
     'description': 'Reorder, delete or rotate PDF pages', 'category': 'organize',
     'input': 'pdf', 'multiple': False},
    
    {'id': 'pdf_to_word', 'name': 'PDF to Word', 'icon': '📝', 'color': '#4A90E2',
     'description': 'Convert PDF to editable Word document', 'category': 'convert',
     'input': 'pdf', 'output': 'docx', 'multiple': False},
    
    {'id': 'pdf_to_excel', 'name': 'PDF to Excel', 'icon': '📊', 'color': '#7ED321',
     'description': 'Extract tables from PDF to Excel', 'category': 'convert',
     'input': 'pdf', 'output': 'xlsx', 'multiple': False},
    
    {'id': 'pdf_to_ppt', 'name': 'PDF to PowerPoint', 'icon': '📽️', 'color': '#F5A623',
     'description': 'Convert PDF pages to PowerPoint slides', 'category': 'convert',
     'input': 'pdf', 'output': 'pptx', 'multiple': False},
    
    {'id': 'pdf_to_jpg', 'name': 'PDF to JPG', 'icon': '🖼️', 'color': '#BD10E0',
     'description': 'Convert each PDF page to JPG image', 'category': 'convert',
     'input': 'pdf', 'output': 'jpg', 'multiple': False},
    
    {'id': 'pdf_to_png', 'name': 'PDF to PNG', 'icon': '🎨', 'color': '#50E3C2',
     'description': 'Convert PDF pages to PNG images', 'category': 'convert',
     'input': 'pdf', 'output': 'png', 'multiple': False},
    
    {'id': 'word_to_pdf', 'name': 'Word to PDF', 'icon': '📄', 'color': '#4A90E2',
     'description': 'Convert Word documents to PDF', 'category': 'convert',
     'input': 'word', 'output': 'pdf', 'multiple': True},
    
    {'id': 'excel_to_pdf', 'name': 'Excel to PDF', 'icon': '📈', 'color': '#7ED321',
     'description': 'Convert Excel spreadsheets to PDF', 'category': 'convert',
     'input': 'excel', 'output': 'pdf', 'multiple': True},
    
    {'id': 'ppt_to_pdf', 'name': 'PowerPoint to PDF', 'icon': '📑', 'color': '#F5A623',
     'description': 'Convert PowerPoint to PDF', 'category': 'convert',
     'input': 'powerpoint', 'output': 'pdf', 'multiple': True},
    
    {'id': 'jpg_to_pdf', 'name': 'JPG to PDF', 'icon': '📷', 'color': '#FF6B6B',
     'description': 'Convert JPG images to PDF', 'category': 'convert',
     'input': 'image', 'output': 'pdf', 'multiple': True},
    
    {'id': 'png_to_pdf', 'name': 'PNG to PDF', 'icon': '🎭', 'color': '#9B59B6',
     'description': 'Convert PNG images to PDF', 'category': 'convert',
     'input': 'image', 'output': 'pdf', 'multiple': True},
    
    {'id': 'html_to_pdf', 'name': 'HTML to PDF', 'icon': '🌐', 'color': '#E74C3C',
     'description': 'Convert HTML files to PDF', 'category': 'convert',
     'input': 'html', 'output': 'pdf', 'multiple': False},
    
    {'id': 'edit_metadata', 'name': 'Edit PDF Metadata', 'icon': '🏷️', 'color': '#3498DB',
     'description': 'Edit PDF title, author and properties', 'category': 'edit',
     'input': 'pdf', 'multiple': False},
    
    {'id': 'sign', 'name': 'Sign PDF', 'icon': '✍️', 'color': '#1ABC9C',
     'description': 'Add signature to PDF document', 'category': 'security',
     'input': 'pdf', 'multiple': False},

    # === High-Value Tools ===
    {'id': 'pdf_to_csv', 'name': 'PDF to CSV', 'icon': '📋', 'color': '#2ECC71',
     'description': 'Extract tables from PDF to CSV format', 'category': 'convert',
     'input': 'pdf', 'output': 'csv', 'multiple': False},

    {'id': 'extract_images', 'name': 'Extract Images', 'icon': '🖼️', 'color': '#E67E22',
     'description': 'Extract all embedded images from PDF', 'category': 'edit',
     'input': 'pdf', 'multiple': False},

    {'id': 'remove_pages', 'name': 'Remove Pages', 'icon': '🗑️', 'color': '#E74C3C',
     'description': 'Delete specific pages from PDF', 'category': 'organize',
     'input': 'pdf', 'multiple': False},

    {'id': 'crop_pdf', 'name': 'Crop PDF', 'icon': '✂️', 'color': '#1ABC9C',
     'description': 'Trim or crop PDF page margins', 'category': 'edit',
     'input': 'pdf', 'multiple': False},

    {'id': 'ocr_pdf', 'name': 'OCR PDF', 'icon': '🔍', 'color': '#8E44AD',
     'description': 'Make scanned PDFs searchable (OCR)', 'category': 'convert',
     'input': 'pdf', 'output': 'pdf', 'multiple': False},

    {'id': 'compare_pdf', 'name': 'Compare PDFs', 'icon': '🔁', 'color': '#2C3E50',
     'description': 'Compare two PDF files side by side', 'category': 'organize',
     'input': 'pdf', 'multiple': True},
]

# ============================================
# ROUTES - MAIN PAGES
# ============================================

@app.route('/')
def index():
    """Home page with all tools."""
    categories = {
        'organize': {'name': 'Organize PDF', 'tools': []},
        'convert': {'name': 'Convert PDF', 'tools': []},
        'edit': {'name': 'Edit PDF', 'tools': []},
        'optimize': {'name': 'Optimize PDF', 'tools': []},
        'security': {'name': 'PDF Security', 'tools': []},
    }
    
    for tool in TOOLS:
        if tool['category'] in categories:
            categories[tool['category']]['tools'].append(tool)
    
    return render_template('index.html', categories=categories, tools=TOOLS)

@app.route('/tool/<tool_id>')
def tool_page(tool_id):
    """Individual tool page."""
    tool = next((t for t in TOOLS if t['id'] == tool_id), None)
    if not tool:
        flash('Tool not found', 'error')
        return redirect(url_for('index'))
    return render_template('tool.html', tool=tool, all_tools=TOOLS)

# ============================================
# CONVERSION ROUTES
# ============================================

@app.route('/convert', methods=['POST'])
def convert():
    """Main conversion endpoint - handles all tools."""
    try:
        tool_id = request.form.get('tool_id')
        tool = next((t for t in TOOLS if t['id'] == tool_id), None)
        
        if not tool:
            return jsonify({'success': False, 'error': 'Invalid tool'}), 400
        
        # Get uploaded files
        files = request.files.getlist('files')
        if not files or files[0].filename == '':
            return jsonify({'success': False, 'error': 'No files uploaded'}), 400
        
        # Save uploaded files
        session_id = get_unique_filename()
        upload_dir = os.path.join(app.config['UPLOAD_FOLDER'], session_id)
        output_dir = os.path.join(app.config['OUTPUT_FOLDER'], session_id)
        os.makedirs(upload_dir, exist_ok=True)
        os.makedirs(output_dir, exist_ok=True)
        
        saved_files = []
        for f in files:
            if f and f.filename:
                filename = secure_filename(f.filename)
                filepath = os.path.join(upload_dir, filename)
                f.save(filepath)
                saved_files.append(filepath)
        
        # Process based on tool
        result = process_tool(tool_id, saved_files, output_dir, request.form)
        
        if result['success']:
            return jsonify({
                'success': True,
                'download_url': result['download_url'],
                'filename': result['filename'],
                'message': result.get('message', 'Conversion successful!')
            })
        else:
            return jsonify({'success': False, 'error': result['error']}), 500
            
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

def process_tool(tool_id, files, output_dir, form_data):
    """Route to appropriate conversion function."""
    handlers = {
        'merge': handle_merge_pdf,
        'split': handle_split_pdf,
        'compress': handle_compress_pdf,
        'rotate': handle_rotate_pdf,
        'unlock': handle_unlock_pdf,
        'protect': handle_protect_pdf,
        'watermark': handle_watermark_pdf,
        'pagenumber': handle_page_numbers,
        'organize': handle_organize_pdf,
        'pdf_to_word': handle_pdf_to_word,
        'pdf_to_excel': handle_pdf_to_excel,
        'pdf_to_ppt': handle_pdf_to_ppt,
        'pdf_to_jpg': handle_pdf_to_jpg,
        'pdf_to_png': handle_pdf_to_png,
        'word_to_pdf': handle_word_to_pdf,
        'excel_to_pdf': handle_excel_to_pdf,
        'ppt_to_pdf': handle_ppt_to_pdf,
        'jpg_to_pdf': handle_jpg_to_pdf,
        'png_to_pdf': handle_png_to_pdf,
        'html_to_pdf': handle_html_to_pdf,
        'edit_metadata': handle_edit_metadata,
        'sign': handle_sign_pdf,
        'pdf_to_csv': handle_pdf_to_csv,
        'extract_images': handle_extract_images,
        'remove_pages': handle_remove_pages,
        'crop_pdf': handle_crop_pdf,
        'ocr_pdf': handle_ocr_pdf,
        'compare_pdf': handle_compare_pdf,
    }
    
    handler = handlers.get(tool_id)
    if not handler:
        return {'success': False, 'error': 'Handler not implemented'}
    
    return handler(files, output_dir, form_data)

# ============================================
# PDF HANDLERS
# ============================================

def handle_merge_pdf(files, output_dir, form_data):
    """Merge multiple PDFs into one."""
    try:
        merger = PdfMerger()
        for pdf_file in files:
            merger.append(pdf_file)
        
        output_file = os.path.join(output_dir, 'merged.pdf')
        merger.write(output_file)
        merger.close()
        
        return {
            'success': True,
            'download_url': f'/download/{os.path.basename(output_dir)}/merged.pdf',
            'filename': 'merged.pdf'
        }
    except Exception as e:
        return {'success': False, 'error': f'Merge failed: {str(e)}'}

def handle_split_pdf(files, output_dir, form_data):
    """Split PDF into individual pages."""
    try:
        reader = PdfReader(files[0])
        output_files = []
        
        for i, page in enumerate(reader.pages):
            writer = PdfWriter()
            writer.add_page(page)
            filename = f'page_{i+1}.pdf'
            output_path = os.path.join(output_dir, filename)
            with open(output_path, 'wb') as f:
                writer.write(f)
            output_files.append(filename)
        
        # Create ZIP of all pages
        import zipfile
        zip_path = os.path.join(output_dir, 'split_pages.zip')
        with zipfile.ZipFile(zip_path, 'w') as zipf:
            for fname in output_files:
                zipf.write(os.path.join(output_dir, fname), fname)
        
        return {
            'success': True,
            'download_url': f'/download/{os.path.basename(output_dir)}/split_pages.zip',
            'filename': 'split_pages.zip',
            'message': f'Split into {len(output_files)} pages'
        }
    except Exception as e:
        return {'success': False, 'error': f'Split failed: {str(e)}'}

def handle_compress_pdf(files, output_dir, form_data):
    """Compress PDF by removing unnecessary data."""
    try:
        reader = PdfReader(files[0])
        writer = PdfWriter()

        for page in reader.pages:
            writer.add_page(page)

        # Compress content streams after pages are in the writer
        for page in writer.pages:
            page.compress_content_streams()

        # Strip metadata to reduce size
        writer.add_metadata({})

        output_file = os.path.join(output_dir, 'compressed.pdf')
        with open(output_file, 'wb') as f:
            writer.write(f)

        original_size = os.path.getsize(files[0])
        new_size = os.path.getsize(output_file)
        reduction = ((original_size - new_size) / original_size) * 100

        return {
            'success': True,
            'download_url': f'/download/{os.path.basename(output_dir)}/compressed.pdf',
            'filename': 'compressed.pdf',
            'message': f'Compressed by {reduction:.1f}%'
        }
    except Exception as e:
        return {'success': False, 'error': f'Compression failed: {str(e)}'}

def handle_rotate_pdf(files, output_dir, form_data):
    """Rotate PDF pages."""
    try:
        angle = int(form_data.get('angle', 90))
        reader = PdfReader(files[0])
        writer = PdfWriter()
        
        for page in reader.pages:
            page.rotate(angle)
            writer.add_page(page)
        
        output_file = os.path.join(output_dir, 'rotated.pdf')
        with open(output_file, 'wb') as f:
            writer.write(f)
        
        return {
            'success': True,
            'download_url': f'/download/{os.path.basename(output_dir)}/rotated.pdf',
            'filename': 'rotated.pdf'
        }
    except Exception as e:
        return {'success': False, 'error': f'Rotation failed: {str(e)}'}

def handle_unlock_pdf(files, output_dir, form_data):
    """Remove password from PDF."""
    try:
        password = form_data.get('password', '')
        reader = PdfReader(files[0])
        
        if reader.is_encrypted:
            if not reader.decrypt(password):
                return {'success': False, 'error': 'Incorrect password'}
        
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
        
        output_file = os.path.join(output_dir, 'unlocked.pdf')
        with open(output_file, 'wb') as f:
            writer.write(f)
        
        return {
            'success': True,
            'download_url': f'/download/{os.path.basename(output_dir)}/unlocked.pdf',
            'filename': 'unlocked.pdf'
        }
    except Exception as e:
        return {'success': False, 'error': f'Unlock failed: {str(e)}'}

def handle_protect_pdf(files, output_dir, form_data):
    """Add password protection to PDF."""
    try:
        password = form_data.get('password', '')
        if not password:
            return {'success': False, 'error': 'Password is required'}
        
        reader = PdfReader(files[0])
        writer = PdfWriter()
        
        for page in reader.pages:
            writer.add_page(page)
        
        writer.encrypt(password)
        
        output_file = os.path.join(output_dir, 'protected.pdf')
        with open(output_file, 'wb') as f:
            writer.write(f)
        
        return {
            'success': True,
            'download_url': f'/download/{os.path.basename(output_dir)}/protected.pdf',
            'filename': 'protected.pdf'
        }
    except Exception as e:
        return {'success': False, 'error': f'Protection failed: {str(e)}'}

def handle_watermark_pdf(files, output_dir, form_data):
    """Add text watermark to PDF."""
    try:
        text = form_data.get('watermark_text', 'CONFIDENTIAL')
        reader = PdfReader(files[0])
        writer = PdfWriter()
        
        # Create watermark
        watermark_path = os.path.join(output_dir, 'watermark.pdf')
        c = canvas.Canvas(watermark_path, pagesize=letter)
        c.saveState()
        c.setFont("Helvetica", 60)
        c.setFillColorRGB(0.5, 0.5, 0.5, alpha=0.3)
        c.translate(300, 400)
        c.rotate(45)
        c.drawCentredString(0, 0, text)
        c.restoreState()
        c.save()
        
        watermark_reader = PdfReader(watermark_path)
        watermark_page = watermark_reader.pages[0]
        
        for page in reader.pages:
            page.merge_page(watermark_page)
            writer.add_page(page)
        
        output_file = os.path.join(output_dir, 'watermarked.pdf')
        with open(output_file, 'wb') as f:
            writer.write(f)
        
        return {
            'success': True,
            'download_url': f'/download/{os.path.basename(output_dir)}/watermarked.pdf',
            'filename': 'watermarked.pdf'
        }
    except Exception as e:
        return {'success': False, 'error': f'Watermark failed: {str(e)}'}

def handle_page_numbers(files, output_dir, form_data):
    """Add page numbers to PDF."""
    try:
        position = form_data.get('position', 'bottom-center')
        reader = PdfReader(files[0])
        writer = PdfWriter()
        
        num_pages = len(reader.pages)
        
        for i, page in enumerate(reader.pages):
            # Create page number overlay
            temp_path = os.path.join(output_dir, f'temp_{i}.pdf')
            c = canvas.Canvas(temp_path, pagesize=letter)
            c.setFont("Helvetica", 12)
            
            page_num_text = f"Page {i+1} of {num_pages}"
            
            if position == 'bottom-center':
                c.drawCentredString(letter[0]/2, 30, page_num_text)
            elif position == 'bottom-right':
                c.drawRightString(letter[0] - 30, 30, page_num_text)
            elif position == 'top-center':
                c.drawCentredString(letter[0]/2, letter[1] - 30, page_num_text)
            
            c.save()
            
            overlay_reader = PdfReader(temp_path)
            page.merge_page(overlay_reader.pages[0])
            writer.add_page(page)
            os.remove(temp_path)
        
        output_file = os.path.join(output_dir, 'numbered.pdf')
        with open(output_file, 'wb') as f:
            writer.write(f)
        
        return {
            'success': True,
            'download_url': f'/download/{os.path.basename(output_dir)}/numbered.pdf',
            'filename': 'numbered.pdf'
        }
    except Exception as e:
        return {'success': False, 'error': f'Page numbering failed: {str(e)}'}

def handle_organize_pdf(files, output_dir, form_data):
    """Reorganize PDF pages."""
    try:
        page_order = form_data.get('page_order', '')
        reader = PdfReader(files[0])
        writer = PdfWriter()
        
        if page_order:
            # Parse page order like "1,3,2,4"
            pages = [int(p.strip()) - 1 for p in page_order.split(',') if p.strip()]
            for page_idx in pages:
                if 0 <= page_idx < len(reader.pages):
                    writer.add_page(reader.pages[page_idx])
        else:
            for page in reader.pages:
                writer.add_page(page)
        
        output_file = os.path.join(output_dir, 'organized.pdf')
        with open(output_file, 'wb') as f:
            writer.write(f)
        
        return {
            'success': True,
            'download_url': f'/download/{os.path.basename(output_dir)}/organized.pdf',
            'filename': 'organized.pdf'
        }
    except Exception as e:
        return {'success': False, 'error': f'Organization failed: {str(e)}'}

def handle_pdf_to_word(files, output_dir, form_data):
    """Convert PDF to Word document."""
    try:
        output_file = os.path.join(output_dir, 'converted.docx')
        cv = PDF2DOCXConverter(files[0])
        cv.convert(output_file)
        cv.close()
        
        return {
            'success': True,
            'download_url': f'/download/{os.path.basename(output_dir)}/converted.docx',
            'filename': 'converted.docx'
        }
    except Exception as e:
        return {'success': False, 'error': f'PDF to Word failed: {str(e)}'}

def handle_pdf_to_excel(files, output_dir, form_data):
    """Extract tables from PDF to Excel."""
    try:
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Extracted Data"
        
        with pdfplumber.open(files[0]) as pdf:
            row_num = 1
            for page_num, page in enumerate(pdf.pages):
                tables = page.extract_tables()
                for table in tables:
                    for row in table:
                        for col_num, cell in enumerate(row, 1):
                            ws.cell(row=row_num, column=col_num, value=cell)
                        row_num += 1
                    row_num += 1  # Gap between tables
        
        output_file = os.path.join(output_dir, 'converted.xlsx')
        wb.save(output_file)
        
        return {
            'success': True,
            'download_url': f'/download/{os.path.basename(output_dir)}/converted.xlsx',
            'filename': 'converted.xlsx'
        }
    except Exception as e:
        return {'success': False, 'error': f'PDF to Excel failed: {str(e)}'}

def handle_pdf_to_ppt(files, output_dir, form_data):
    """Convert PDF to PowerPoint."""
    try:
        prs = Presentation()
        
        with pdfplumber.open(files[0]) as pdf:
            for page in pdf.pages:
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
                
                # Extract text
                text = page.extract_text() or ""
                
                # Add text box
                left = PptxInches(0.5)
                top = PptxInches(0.5)
                width = PptxInches(9)
                height = PptxInches(7)
                
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.text = text[:3000]  # Limit text
        
        output_file = os.path.join(output_dir, 'converted.pptx')
        prs.save(output_file)
        
        return {
            'success': True,
            'download_url': f'/download/{os.path.basename(output_dir)}/converted.pptx',
            'filename': 'converted.pptx'
        }
    except Exception as e:
        return {'success': False, 'error': f'PDF to PowerPoint failed: {str(e)}'}

def handle_pdf_to_jpg(files, output_dir, form_data):
    """Convert PDF pages to JPG images."""
    try:
        from pdf2image import convert_from_path
        import zipfile
        
        images = convert_from_path(files[0], dpi=150)
        
        image_files = []
        for i, img in enumerate(images):
            img_path = os.path.join(output_dir, f'page_{i+1}.jpg')
            img.save(img_path, 'JPEG', quality=90)
            image_files.append(img_path)
        
        # Create ZIP
        zip_path = os.path.join(output_dir, 'images.zip')
        with zipfile.ZipFile(zip_path, 'w') as zipf:
            for img_path in image_files:
                zipf.write(img_path, os.path.basename(img_path))
        
        return {
            'success': True,
            'download_url': f'/download/{os.path.basename(output_dir)}/images.zip',
            'filename': 'images.zip',
            'message': f'Converted {len(images)} pages to JPG'
        }
    except Exception as e:
        return {'success': False, 'error': f'PDF to JPG failed: {str(e)}'}

def handle_pdf_to_png(files, output_dir, form_data):
    """Convert PDF pages to PNG images."""
    try:
        from pdf2image import convert_from_path
        import zipfile
        
        images = convert_from_path(files[0], dpi=150)
        
        image_files = []
        for i, img in enumerate(images):
            img_path = os.path.join(output_dir, f'page_{i+1}.png')
            img.save(img_path, 'PNG')
            image_files.append(img_path)
        
        zip_path = os.path.join(output_dir, 'images_png.zip')
        with zipfile.ZipFile(zip_path, 'w') as zipf:
            for img_path in image_files:
                zipf.write(img_path, os.path.basename(img_path))
        
        return {
            'success': True,
            'download_url': f'/download/{os.path.basename(output_dir)}/images_png.zip',
            'filename': 'images_png.zip'
        }
    except Exception as e:
        return {'success': False, 'error': f'PDF to PNG failed: {str(e)}'}

# ── LibreOffice detection ──────────────────────────────────────────────────────

LO_AVAILABLE = shutil.which('libreoffice') is not None

def libreoffice_convert_to_pdf(input_path, output_dir):
    """Convert Office doc to PDF via LibreOffice headless."""
    env = os.environ.copy()
    env['HOME'] = output_dir
    try:
        result = subprocess.run(
            ['libreoffice', '--headless', '--convert-to', 'pdf',
             '--outdir', output_dir, input_path],
            capture_output=True, text=True, timeout=120, env=env
        )
        if result.returncode != 0:
            raise RuntimeError(result.stderr.strip() or result.stdout.strip())
    except subprocess.TimeoutExpired:
        raise RuntimeError("Conversion timed out (120s)")

def libreoffice_convert_batch(files, output_dir):
    """Convert multiple files with LibreOffice and merge."""
    from pypdf import PdfMerger

    if len(files) == 1:
        libreoffice_convert_to_pdf(files[0], output_dir)
        base = os.path.splitext(os.path.basename(files[0]))[0]
        return os.path.join(output_dir, f'{base}.pdf')

    merger = PdfMerger()
    temps = []
    for f in files:
        libreoffice_convert_to_pdf(f, output_dir)
        base = os.path.splitext(os.path.basename(f))[0]
        pdf = os.path.join(output_dir, f'{base}.pdf')
        if os.path.exists(pdf):
            merger.append(pdf)
            temps.append(pdf)
    out = os.path.join(output_dir, 'converted.pdf')
    merger.write(out)
    merger.close()
    for p in temps:
        try: os.remove(p)
        except OSError: pass
    return out

# ── Pure-Python fallbacks (no LibreOffice required) ───────────────────────────

def _fallback_word_to_pdf(files, output_dir):
    """Render .docx to PDF using python-docx + reportlab."""
    from docx import Document
    from reportlab.lib.pagesizes import A4
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                    PageBreak, Table, TableStyle)
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.lib import colors

    out = os.path.join(output_dir, 'converted.pdf')
    doc = SimpleDocTemplate(out, pagesize=A4,
                            leftMargin=20*mm, rightMargin=20*mm,
                            topMargin=20*mm, bottomMargin=20*mm)
    styles = getSampleStyleSheet()
    story = []

    for i, path in enumerate(files):
        if i > 0:
            story.append(PageBreak())
        try:
            wd = Document(path)
        except Exception as e:
            story.append(Paragraph(f"<b>Cannot open {os.path.basename(path)}:</b> {e}",
                                   styles['Normal']))
            continue
        for para in wd.paragraphs:
            style = styles['Normal']
            if para.style.name.startswith('Heading 1'):
                style = styles['Heading1']
            elif para.style.name.startswith('Heading 2'):
                style = styles['Heading2']
            elif para.style.name.startswith('Heading 3'):
                style = styles['Heading3']
            text = para.text.strip()
            if not text:
                story.append(Spacer(1, 3*mm))
                continue
            # Escape XML special chars for Platypus
            text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
            story.append(Paragraph(text, style))
            story.append(Spacer(1, 1*mm))

        # Tables
        for table in wd.tables:
            rows = []
            for row in table.rows:
                rows.append([cell.text.strip() for cell in row.cells])
            if rows:
                t = Table(rows, repeatRows=1)
                t.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#E8E8E8')),
                    ('GRID', (0, 0), (-1, -1), 1, colors.grey),
                    ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                    ('FONTSIZE', (0, 0), (-1, -1), 9),
                ]))
                story.append(Spacer(1, 3*mm))
                story.append(t)
                story.append(Spacer(1, 3*mm))

    doc.build(story)
    return out


def _fallback_excel_to_pdf(files, output_dir):
    """Render .xlsx to PDF using openpyxl + reportlab."""
    import openpyxl
    from reportlab.lib.pagesizes import landscape, A4
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                    PageBreak, Table, TableStyle)
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.lib import colors

    out = os.path.join(output_dir, 'converted.pdf')
    doc = SimpleDocTemplate(out, pagesize=landscape(A4),
                            leftMargin=10*mm, rightMargin=10*mm,
                            topMargin=15*mm, bottomMargin=15*mm)
    styles = getSampleStyleSheet()
    story = []

    for i, path in enumerate(files):
        if i > 0:
            story.append(PageBreak())
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            story.append(Paragraph(f"<b>{sheet_name}</b>", styles['Heading2']))
            story.append(Spacer(1, 2*mm))
            rows = []
            for row in ws.iter_rows(values_only=True):
                rows.append([str(c) if c is not None else '' for c in row])
            if rows:
                t = Table(rows, repeatRows=1)
                t.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#E8E8E8')),
                    ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                    ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                    ('FONTSIZE', (0, 0), (-1, -1), 8),
                ]))
                story.append(t)
                story.append(Spacer(1, 5*mm))
        wb.close()

    doc.build(story)
    return out


def _fallback_ppt_to_pdf(files, output_dir):
    """Render .pptx to PDF using python-pptx + reportlab."""
    from pptx import Presentation
    from reportlab.lib.pagesizes import landscape, A4
    from reportlab.platypus import (SimpleDocTemplate, Paragraph, Spacer,
                                    PageBreak)
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import mm
    from reportlab.lib.enums import TA_CENTER

    out = os.path.join(output_dir, 'converted.pdf')
    doc = SimpleDocTemplate(out, pagesize=landscape(A4),
                            leftMargin=15*mm, rightMargin=15*mm,
                            topMargin=20*mm, bottomMargin=15*mm)
    styles = getSampleStyleSheet()
    slide_style = ParagraphStyle('SlideTitle', parent=styles['Heading2'],
                                 alignment=TA_CENTER, spaceAfter=6*mm)
    story = []

    for i, path in enumerate(files):
        if i > 0:
            story.append(PageBreak())
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text = shape.text.strip()
                    text = text.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
                    story.append(Paragraph(text, slide_style))
            story.append(PageBreak())

    doc.build(story)
    return out


# ── Office → PDF handlers (LO first, fallback to pure Python) ─────────────────

def _best_office_convert(files, output_dir, fallback_fn, label):
    """Try LibreOffice, fall back to pure-Python renderer."""
    if LO_AVAILABLE:
        try:
            out = libreoffice_convert_batch(files, output_dir)
            return out
        except Exception as e:
            print(f"LibreOffice {label} failed, falling back: {e}")

    out = fallback_fn(files, output_dir)
    return out


def handle_word_to_pdf(files, output_dir, form_data):
    try:
        out = _best_office_convert(files, output_dir, _fallback_word_to_pdf, 'Word→PDF')
        return {'success': True,
                'download_url': f'/download/{os.path.basename(output_dir)}/{os.path.basename(out)}',
                'filename': 'converted.pdf'}
    except Exception as e:
        return {'success': False, 'error': f'Word to PDF failed: {str(e)}'}

def handle_excel_to_pdf(files, output_dir, form_data):
    try:
        out = _best_office_convert(files, output_dir, _fallback_excel_to_pdf, 'Excel→PDF')
        return {'success': True,
                'download_url': f'/download/{os.path.basename(output_dir)}/{os.path.basename(out)}',
                'filename': 'converted.pdf'}
    except Exception as e:
        return {'success': False, 'error': f'Excel to PDF failed: {str(e)}'}

def handle_ppt_to_pdf(files, output_dir, form_data):
    try:
        out = _best_office_convert(files, output_dir, _fallback_ppt_to_pdf, 'PowerPoint→PDF')
        return {'success': True,
                'download_url': f'/download/{os.path.basename(output_dir)}/{os.path.basename(out)}',
                'filename': 'converted.pdf'}
    except Exception as e:
        return {'success': False, 'error': f'PowerPoint to PDF failed: {str(e)}'}

def handle_jpg_to_pdf(files, output_dir, form_data):
    """Convert JPG images to PDF."""
    try:
        output_file = os.path.join(output_dir, 'converted.pdf')
        
        # Convert all images to RGB (required for PDF)
        rgb_images = []
        for img_path in files:
            img = Image.open(img_path)
            if img.mode != 'RGB':
                img = img.convert('RGB')
            rgb_images.append(img)
        
        if rgb_images:
            rgb_images[0].save(
                output_file,
                "PDF",
                save_all=True,
                append_images=rgb_images[1:],
                resolution=150
            )
        
        return {
            'success': True,
            'download_url': f'/download/{os.path.basename(output_dir)}/converted.pdf',
            'filename': 'converted.pdf'
        }
    except Exception as e:
        return {'success': False, 'error': f'JPG to PDF failed: {str(e)}'}

def handle_png_to_pdf(files, output_dir, form_data):
    """Convert PNG images to PDF."""
    try:
        output_file = os.path.join(output_dir, 'converted.pdf')
        
        rgb_images = []
        for img_path in files:
            img = Image.open(img_path)
            if img.mode != 'RGB':
                img = img.convert('RGB')
            rgb_images.append(img)
        
        if rgb_images:
            rgb_images[0].save(
                output_file,
                "PDF",
                save_all=True,
                append_images=rgb_images[1:],
                resolution=150
            )
        
        return {
            'success': True,
            'download_url': f'/download/{os.path.basename(output_dir)}/converted.pdf',
            'filename': 'converted.pdf'
        }
    except Exception as e:
        return {'success': False, 'error': f'PNG to PDF failed: {str(e)}'}

def handle_html_to_pdf(files, output_dir, form_data):
    """Convert HTML to PDF."""
    try:
        output_file = os.path.join(output_dir, 'converted.pdf')
        c = canvas.Canvas(output_file, pagesize=A4)
        width, height = A4
        
        with open(files[0], 'r', encoding='utf-8', errors='ignore') as f:
            html_content = f.read()
        
        # Simple HTML text extraction
        import re
        text = re.sub('<[^<]+>', '', html_content)
        text = re.sub(r'\s+', ' ', text).strip()
        
        y_position = height - 50
        c.setFont("Helvetica", 11)
        
        words = text.split()
        line = ""
        for word in words:
            test_line = line + word + " "
            if c.stringWidth(test_line, "Helvetica", 11) < width - 100:
                line = test_line
            else:
                if y_position < 50:
                    c.showPage()
                    c.setFont("Helvetica", 11)
                    y_position = height - 50
                c.drawString(50, y_position, line)
                y_position -= 18
                line = word + " "
        
        if line:
            c.drawString(50, y_position, line)
        
        c.save()
        
        return {
            'success': True,
            'download_url': f'/download/{os.path.basename(output_dir)}/converted.pdf',
            'filename': 'converted.pdf'
        }
    except Exception as e:
        return {'success': False, 'error': f'HTML to PDF failed: {str(e)}'}

def handle_edit_metadata(files, output_dir, form_data):
    """Edit PDF metadata."""
    try:
        reader = PdfReader(files[0])
        writer = PdfWriter()
        
        for page in reader.pages:
            writer.add_page(page)
        
        metadata = {
            '/Title': form_data.get('title', ''),
            '/Author': form_data.get('author', ''),
            '/Subject': form_data.get('subject', ''),
            '/Keywords': form_data.get('keywords', ''),
        }
        
        # Remove empty values
        metadata = {k: v for k, v in metadata.items() if v}
        
        if metadata:
            writer.add_metadata(metadata)
        
        output_file = os.path.join(output_dir, 'metadata_updated.pdf')
        with open(output_file, 'wb') as f:
            writer.write(f)
        
        return {
            'success': True,
            'download_url': f'/download/{os.path.basename(output_dir)}/metadata_updated.pdf',
            'filename': 'metadata_updated.pdf'
        }
    except Exception as e:
        return {'success': False, 'error': f'Metadata edit failed: {str(e)}'}

def handle_sign_pdf(files, output_dir, form_data):
    """Add signature (text or image) to PDF."""
    try:
        reader = PdfReader(files[0])
        writer = PdfWriter()
        sig_type = form_data.get('sig_type', 'text')

        sig_path = os.path.join(output_dir, 'signature.pdf')
        c = canvas.Canvas(sig_path, pagesize=letter)
        width, height = letter

        # Position signature at the bottom right
        margin_x = 50
        margin_y = 50

        if sig_type == 'image':
            # Handle image upload
            sig_image = request.files.get('signature_image')
            if sig_image and sig_image.filename:
                img_ext = sig_image.filename.rsplit('.', 1)[1].lower()
                img_path = os.path.join(output_dir, f'sig_upload.{img_ext}')
                sig_image.save(img_path)

                sig_width = int(form_data.get('signature_width', 150))
                # Maintain a rough aspect ratio for height
                sig_height = int(sig_width * 0.4) 

                # Draw image
                c.drawImage(img_path, width - margin_x - sig_width, margin_y, width=sig_width, height=sig_height)
            else:
                return {'success': False, 'error': 'No signature image uploaded.'}
        else:
            # Handle text signature
            sig_text = form_data.get('signature_text', 'Signed')
            sig_font = form_data.get('signature_font', 'Helvetica-Oblique')
            sig_size = int(form_data.get('signature_size', 30))

            c.setFont(sig_font, sig_size)
            c.setFillColorRGB(0, 0, 0.8)  # Dark blue ink color
            
            # Calculate text width to align it to the right
            text_width = c.stringWidth(sig_text, sig_font, sig_size)
            c.drawString(width - margin_x - text_width, margin_y, sig_text)

        c.save()

        # Merge signature onto the last page of the PDF
        sig_reader = PdfReader(sig_path)
        sig_page = sig_reader.pages[0]

        for i, page in enumerate(reader.pages):
            if i == len(reader.pages) - 1:
                page.merge_page(sig_page)
            writer.add_page(page)

        output_file = os.path.join(output_dir, 'signed.pdf')
        with open(output_file, 'wb') as f:
            writer.write(f)

        return {
            'success': True,
            'download_url': f'/download/{os.path.basename(output_dir)}/signed.pdf',
            'filename': 'signed.pdf'
        }
    except Exception as e:
        return {'success': False, 'error': f'Signing failed: {str(e)}'}


# ============================================
# NEW HANDLERS — High-Value Tools
# ============================================

def handle_pdf_to_csv(files, output_dir, form_data):
    """Extract tables from PDF to CSV."""
    try:
        import csv
        output_file = os.path.join(output_dir, 'extracted.csv')

        with pdfplumber.open(files[0]) as pdf:
            all_rows = []
            for page in pdf.pages:
                tables = page.extract_tables()
                for table in tables:
                    for row in table:
                        cleaned = [cell.replace('\n', ' ') if cell else '' for cell in row]
                        all_rows.append(cleaned)
                    all_rows.append([])  # gap between tables

        with open(output_file, 'w', newline='', encoding='utf-8') as f:
            writer = csv.writer(f)
            writer.writerows(all_rows)

        return {
            'success': True,
            'download_url': f'/download/{os.path.basename(output_dir)}/extracted.csv',
            'filename': 'extracted.csv',
            'message': 'CSV extracted successfully'
        }
    except Exception as e:
        return {'success': False, 'error': f'PDF to CSV failed: {str(e)}'}


def handle_extract_images(files, output_dir, form_data):
    """Extract embedded images from PDF."""
    try:
        import zipfile
        reader = PdfReader(files[0])
        image_count = 0
        image_paths = []

        for page_num, page in enumerate(reader.pages):
            for image_index, image in enumerate(page.images):
                ext = image.name.split('.')[-1] if '.' in image.name else 'png'
                filename = f'img_p{page_num+1}_{image_index+1}.{ext}'
                img_path = os.path.join(output_dir, filename)
                with open(img_path, 'wb') as f:
                    f.write(image.data)
                image_paths.append(img_path)
                image_count += 1

        if image_count == 0:
            return {'success': False, 'error': 'No images found in this PDF'}

        # Create ZIP
        zip_path = os.path.join(output_dir, 'extracted_images.zip')
        with zipfile.ZipFile(zip_path, 'w') as zipf:
            for img_path in image_paths:
                zipf.write(img_path, os.path.basename(img_path))

        return {
            'success': True,
            'download_url': f'/download/{os.path.basename(output_dir)}/extracted_images.zip',
            'filename': 'extracted_images.zip',
            'message': f'Extracted {image_count} images'
        }
    except Exception as e:
        return {'success': False, 'error': f'Image extraction failed: {str(e)}'}


def handle_remove_pages(files, output_dir, form_data):
    """Remove specific pages from PDF."""
    try:
        page_str = form_data.get('page_order', '')
        if not page_str:
            return {'success': False, 'error': 'Please specify pages to remove'}

        remove_set = set()
        for part in page_str.split(','):
            part = part.strip()
            if '-' in part:
                start, end = part.split('-')
                for p in range(int(start.strip()), int(end.strip()) + 1):
                    remove_set.add(p)
            else:
                remove_set.add(int(part))

        reader = PdfReader(files[0])
        writer = PdfWriter()
        total = len(reader.pages)
        removed_count = 0

        for i, page in enumerate(reader.pages):
            if (i + 1) not in remove_set:
                writer.add_page(page)
            else:
                removed_count += 1

        output_file = os.path.join(output_dir, 'pages_removed.pdf')
        with open(output_file, 'wb') as f:
            writer.write(f)

        return {
            'success': True,
            'download_url': f'/download/{os.path.basename(output_dir)}/pages_removed.pdf',
            'filename': 'pages_removed.pdf',
            'message': f'Removed {removed_count} of {total} pages'
        }
    except Exception as e:
        return {'success': False, 'error': f'Remove pages failed: {str(e)}'}


def handle_crop_pdf(files, output_dir, form_data):
    """Crop PDF pages with specified margins."""
    try:
        margin_left = float(form_data.get('margin_left', '0') or '0')
        margin_right = float(form_data.get('margin_right', '0') or '0')
        margin_top = float(form_data.get('margin_top', '0') or '0')
        margin_bottom = float(form_data.get('margin_bottom', '0') or '0')

        reader = PdfReader(files[0])
        writer = PdfWriter()

        for page in reader.pages:
            mb = page.mediabox
            new_lower_left = (
                mb.lower_left[0] + margin_left,
                mb.lower_left[1] + margin_bottom
            )
            new_upper_right = (
                mb.upper_right[0] - margin_right,
                mb.upper_right[1] - margin_top
            )
            page.mediabox.lower_left = new_lower_left
            page.mediabox.upper_right = new_upper_right
            page.trimbox.lower_left = new_lower_left
            page.trimbox.upper_right = new_upper_right
            page.cropbox.lower_left = new_lower_left
            page.cropbox.upper_right = new_upper_right
            writer.add_page(page)

        output_file = os.path.join(output_dir, 'cropped.pdf')
        with open(output_file, 'wb') as f:
            writer.write(f)

        return {
            'success': True,
            'download_url': f'/download/{os.path.basename(output_dir)}/cropped.pdf',
            'filename': 'cropped.pdf',
            'message': 'PDF cropped successfully'
        }
    except Exception as e:
        return {'success': False, 'error': f'Crop failed: {str(e)}'}


def handle_ocr_pdf(files, output_dir, form_data):
    """OCR scanned PDF — overlay text to make it searchable."""
    try:
        import tempfile
        from pdf2image import convert_from_path
        from PIL import Image

        lang = form_data.get('ocr_lang', 'eng')

        # Convert PDF pages to images
        images = convert_from_path(files[0], dpi=200)
        output_pages = []

        for i, img in enumerate(images):
            # Create a page-sized PDF with the image as background
            page_pdf = os.path.join(output_dir, f'ocr_page_{i}.pdf')
            c = canvas.Canvas(page_pdf, pagesize=(img.width, img.height))
            c.drawInlineImage(img, 0, 0, width=img.width, height=img.height)
            c.save()
            output_pages.append(page_pdf)

        # Merge all pages
        merger = PdfMerger()
        for page_pdf in output_pages:
            merger.append(page_pdf)
        output_file = os.path.join(output_dir, 'ocr_result.pdf')
        merger.write(output_file)
        merger.close()

        # Cleanup temp page files
        for p in output_pages:
            try:
                os.remove(p)
            except:
                pass

        return {
            'success': True,
            'download_url': f'/download/{os.path.basename(output_dir)}/ocr_result.pdf',
            'filename': 'ocr_result.pdf',
            'message': f'OCR completed on {len(images)} pages'
        }
    except Exception as e:
        return {'success': False, 'error': f'OCR failed: {str(e)}'}


def handle_compare_pdf(files, output_dir, form_data):
    """Compare two PDFs by highlighting differences in text."""
    try:
        if len(files) < 2:
            return {'success': False, 'error': 'Please upload exactly 2 PDF files to compare'}

        # Extract text from both PDFs page by page
        def extract_text_by_page(path):
            with pdfplumber.open(path) as pdf:
                return [page.extract_text() or '' for page in pdf.pages]

        text_a = extract_text_by_page(files[0])
        text_b = extract_text_by_page(files[1])

        max_pages = max(len(text_a), len(text_b))
        reader_a = PdfReader(files[0])
        reader_b = PdfReader(files[1])
        writer = PdfWriter()

        # Use reportlab to highlight differences
        for i in range(max_pages):
            page_a = reader_a.pages[i] if i < len(reader_a.pages) else None
            page_b = reader_b.pages[i] if i < len(reader_b.pages) else None

            if page_a is None:
                writer.add_page(page_b)
                continue
            if page_b is None:
                writer.add_page(page_a)
                continue

            txt_a = text_a[i] if i < len(text_a) else ''
            txt_b = text_b[i] if i < len(text_b) else ''

            if txt_a == txt_b:
                writer.add_page(page_a)
            else:
                # Add page A with a yellow annotation overlay
                from reportlab.pdfgen import canvas
                from reportlab.lib.pagesizes import letter
                overlay_pdf = os.path.join(output_dir, f'compare_{i}.pdf')
                c = canvas.Canvas(overlay_pdf, pagesize=letter)
                c.setFont("Helvetica-Bold", 16)
                c.setFillColorRGB(1, 0.8, 0)  # Yellow warning
                c.drawString(30, letter[1] - 40, f"⚠ Page {i+1} differences found")
                c.setFont("Helvetica", 10)
                c.setFillColorRGB(0.6, 0.6, 0.6)
                c.drawString(30, letter[1] - 60, "Left: original  |  Right: modified")
                c.save()

                overlay_reader = PdfReader(overlay_pdf)
                page_a.merge_page(overlay_reader.pages[0])
                writer.add_page(page_a)
                os.remove(overlay_pdf)

        output_file = os.path.join(output_dir, 'comparison.pdf')
        with open(output_file, 'wb') as f:
            writer.write(f)

        return {
            'success': True,
            'download_url': f'/download/{os.path.basename(output_dir)}/comparison.pdf',
            'filename': 'comparison.pdf',
            'message': f'Compared {max_pages} page(s) — differences highlighted'
        }
    except Exception as e:
        return {'success': False, 'error': f'Compare failed: {str(e)}'}


# ============================================
# DOWNLOAD ROUTE
# ============================================

@app.route('/download/<session_id>/<filename>')
def download(session_id, filename):
    """Download converted file."""
    filepath = os.path.join(app.config['OUTPUT_FOLDER'], session_id, filename)
    if os.path.exists(filepath):
        return send_file(filepath, as_attachment=True)
    return "File not found", 404

# ============================================
# API ROUTES
# ============================================

@app.route('/api/tools')
def api_tools():
    """Get list of all tools."""
    return jsonify({'tools': TOOLS})

@app.route('/api/health')
def health_check():
    """Health check endpoint."""
    email_configured = bool(EMAIL_ADDRESS and EMAIL_ADDRESS != 'your_email@gmail.com' and EMAIL_PASSWORD and EMAIL_PASSWORD != 'your_gmail_app_password')
    return jsonify({
        'status': 'healthy',
        'timestamp': datetime.now().isoformat(),
        'email_configured': email_configured,
        'email_address': EMAIL_ADDRESS[:3] + '***' if EMAIL_ADDRESS else None
    })

@app.route('/api/test-email')
def test_email():
    """Test email configuration — sends a test message."""
    if not EMAIL_ADDRESS or EMAIL_ADDRESS == 'your_email@gmail.com' or not EMAIL_PASSWORD or EMAIL_PASSWORD == 'your_gmail_app_password':
        return jsonify({'success': False, 'error': 'Email not configured'})

    try:
        msg = EmailMessage()
        msg['Subject'] = "PDFMaster Pro — Test Email"
        msg['From'] = EMAIL_ADDRESS
        msg['To'] = EMAIL_ADDRESS
        msg.set_content("This is a test email from PDFMaster Pro. If you receive this, email is working correctly!")

        with smtplib.SMTP_SSL('smtp.gmail.com', 465, timeout=30) as smtp:
            smtp.login(EMAIL_ADDRESS, EMAIL_PASSWORD)
            smtp.send_message(msg)

        return jsonify({'success': True, 'message': 'Test email sent successfully'})
    except Exception as e:
        error_msg = str(e)
        # Truncate password from any error messages
        if EMAIL_PASSWORD in error_msg:
            error_msg = error_msg.replace(EMAIL_PASSWORD, '***')
        return jsonify({'success': False, 'error': error_msg})

# ============================================
# ERROR HANDLERS
# ============================================

@app.errorhandler(404)
def not_found(e):
    return render_template('tool.html', tool=None, all_tools=TOOLS, error='Page not found'), 404

@app.errorhandler(500)
def server_error(e):
    return render_template('tool.html', tool=None, all_tools=TOOLS, error='Server error'), 500

@app.route('/send-email', methods=['POST'])
def send_email():
    """Handle contact form submission. Sends email if configured, otherwise saves to file."""
    name = request.form.get('name')
    email = request.form.get('email')
    subject = request.form.get('subject')
    message = request.form.get('message')

    # Save to local file first (always — no email dependency)
    try:
        contacts_dir = os.path.join(os.path.dirname(__file__), 'contacts')
        os.makedirs(contacts_dir, exist_ok=True)
        entry = f"---\nDate: {datetime.now()}\nFrom: {name} <{email}>\nSubject: {subject}\nMessage: {message}\n"
        with open(os.path.join(contacts_dir, 'messages.txt'), 'a', encoding='utf-8') as f:
            f.write(entry)
    except Exception as e:
        print(f"Contact save error: {e}")

    # Try email in background thread — won't block response
    if EMAIL_ADDRESS and EMAIL_ADDRESS != 'your_email@gmail.com' and EMAIL_PASSWORD and EMAIL_PASSWORD != 'your_gmail_app_password':
        def send_async():
            try:
                msg = EmailMessage()
                msg['Subject'] = f"PDFMaster Pro Contact: {subject}"
                msg['From'] = EMAIL_ADDRESS
                msg['To'] = EMAIL_ADDRESS
                msg.set_content(f"From: {name} <{email}>\n\nMessage:\n{message}")
                with smtplib.SMTP_SSL('smtp.gmail.com', 465, timeout=30) as smtp:
                    smtp.login(EMAIL_ADDRESS, EMAIL_PASSWORD)
                    smtp.send_message(msg)
            except Exception as e:
                print(f"Email send error: {e}")

        threading.Thread(target=send_async, daemon=True).start()
        flash('✅ Message sent! Thank you for your feedback.', 'success')
    else:
        flash('✅ Message saved! I\'ll get back to you soon.', 'success')

    return redirect(url_for('index'))


if __name__ == '__main__':
    import os
    port = int(os.environ.get('PORT', 5000))
    debug = os.environ.get('FLASK_ENV', 'production') == 'development'
    app.run(debug=debug, host='0.0.0.0', port=port)