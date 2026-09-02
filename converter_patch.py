"""Targeted production fixes for PDF compression and PDF-to-PowerPoint."""

import os
import shutil
import subprocess


def handle_compress_pdf(files, output_dir, form_data):
    """Compress a PDF with Ghostscript image downsampling, with a safe fallback."""
    try:
        if not files:
            return {'success': False, 'error': 'No PDF uploaded'}

        src = files[0]
        output_file = os.path.join(output_dir, 'compressed.pdf')
        original_size = os.path.getsize(src)
        gs = shutil.which('gs')

        candidates = []
        if gs:
            for preset, name in (('/ebook', 'ebook.pdf'), ('/screen', 'screen.pdf')):
                candidate = os.path.join(output_dir, name)
                cmd = [
                    gs,
                    '-sDEVICE=pdfwrite',
                    '-dCompatibilityLevel=1.4',
                    f'-dPDFSETTINGS={preset}',
                    '-dNOPAUSE',
                    '-dQUIET',
                    '-dBATCH',
                    '-dDetectDuplicateImages=true',
                    '-dCompressFonts=true',
                    '-dSubsetFonts=true',
                    f'-sOutputFile={candidate}',
                    src,
                ]
                result = subprocess.run(cmd, capture_output=True, text=True, timeout=180)
                if result.returncode == 0 and os.path.exists(candidate) and os.path.getsize(candidate) > 0:
                    candidates.append(candidate)

        if candidates:
            best = min(candidates, key=os.path.getsize)
            shutil.copy2(best if os.path.getsize(best) < original_size else src, output_file)
        else:
            from pypdf import PdfReader, PdfWriter
            reader = PdfReader(src)
            writer = PdfWriter()
            for page in reader.pages:
                writer.add_page(page)
            for page in writer.pages:
                try:
                    page.compress_content_streams()
                except Exception:
                    pass
            writer.add_metadata({})
            with open(output_file, 'wb') as fh:
                writer.write(fh)
            if os.path.getsize(output_file) >= original_size:
                shutil.copy2(src, output_file)

        new_size = os.path.getsize(output_file)
        reduction = max(0.0, ((original_size - new_size) / original_size) * 100) if original_size else 0.0
        if reduction < 0.1:
            message = 'This PDF is already well optimized; no meaningful size reduction was possible.'
        else:
            message = f'Compressed by {reduction:.1f}% ({original_size/1048576:.2f} MB → {new_size/1048576:.2f} MB)'

        return {
            'success': True,
            'download_url': f'/download/{os.path.basename(output_dir)}/compressed.pdf',
            'filename': 'compressed.pdf',
            'message': message,
        }
    except subprocess.TimeoutExpired:
        return {'success': False, 'error': 'Compression timed out. Try a smaller PDF.'}
    except Exception as exc:
        return {'success': False, 'error': f'Compression failed: {exc}'}


def _group_pdf_words_into_lines(words, tolerance=3.0):
    if not words:
        return []
    ordered = sorted(words, key=lambda w: (float(w.get('top', 0)), float(w.get('x0', 0))))
    lines = []
    for word in ordered:
        top = float(word.get('top', 0))
        match = None
        for line in reversed(lines[-4:]):
            if abs(line['top'] - top) <= tolerance:
                match = line
                break
        if match is None:
            match = {'top': top, 'words': []}
            lines.append(match)
        match['words'].append(word)

    result = []
    for line in lines:
        ws = sorted(line['words'], key=lambda w: float(w.get('x0', 0)))
        text = ' '.join((w.get('text') or '').strip() for w in ws if (w.get('text') or '').strip())
        if not text:
            continue
        result.append({
            'text': text,
            'x0': min(float(w.get('x0', 0)) for w in ws),
            'x1': max(float(w.get('x1', 0)) for w in ws),
            'top': min(float(w.get('top', 0)) for w in ws),
            'bottom': max(float(w.get('bottom', 0)) for w in ws),
            'size': max(float(w.get('size') or 10) for w in ws),
        })
    return result


def _region_is_near_white(image, line, page_width, page_height):
    """Return True when the text region is on an ordinary light page background."""
    try:
        iw, ih = image.size
        x0 = max(0, int(line['x0'] / page_width * iw) - 2)
        x1 = min(iw, int(line['x1'] / page_width * iw) + 2)
        y0 = max(0, int(line['top'] / page_height * ih) - 2)
        y1 = min(ih, int(line['bottom'] / page_height * ih) + 2)
        if x1 <= x0 or y1 <= y0:
            return False
        region = image.crop((x0, y0, x1, y1)).convert('RGB').resize((1, 1))
        r, g, b = region.getpixel((0, 0))
        # The sample contains the dark text itself, so a threshold around 200
        # still corresponds to a predominantly white/light background.
        return min(r, g, b) >= 200
    except Exception:
        return False


def handle_pdf_to_ppt(files, output_dir, form_data):
    """Convert each PDF page to a faithful slide and reconstruct editable text."""
    temp_images = []
    try:
        if not files:
            return {'success': False, 'error': 'No PDF uploaded'}

        import pdfplumber
        from pdf2image import convert_from_path
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
        from pptx.util import Inches, Pt

        src = files[0]
        with pdfplumber.open(src) as pdf:
            page_count = len(pdf.pages)
            if page_count == 0:
                return {'success': False, 'error': 'The PDF contains no pages'}
            page_meta = []
            for page in pdf.pages:
                try:
                    words = page.extract_words(
                        use_text_flow=True,
                        keep_blank_chars=False,
                        extra_attrs=['size'],
                    )
                except Exception:
                    words = page.extract_words(use_text_flow=True, keep_blank_chars=False)
                page_meta.append({
                    'width': float(page.width),
                    'height': float(page.height),
                    'lines': _group_pdf_words_into_lines(words),
                })

        images = convert_from_path(
            src,
            dpi=150,
            fmt='jpeg',
            jpegopt={'quality': 88, 'progressive': True, 'optimize': True},
            thread_count=1,
        )
        if len(images) != page_count:
            return {
                'success': False,
                'error': f'PDF rendering returned {len(images)} page(s), but the PDF contains {page_count}. No incomplete PPTX was created.'
            }

        prs = Presentation()
        while len(prs.slides):
            rId = prs.slides._sldIdLst[0].rId
            prs.part.drop_rel(rId)
            del prs.slides._sldIdLst[0]

        first = page_meta[0]
        slide_w = Inches(10)
        slide_h = int(slide_w * first['height'] / first['width'])
        prs.slide_width = slide_w
        prs.slide_height = slide_h
        editable_line_count = 0

        for idx, (image, meta) in enumerate(zip(images, page_meta), start=1):
            img_path = os.path.join(output_dir, f'ppt_page_{idx}.jpg')
            image.save(img_path, 'JPEG', quality=88, optimize=True)
            temp_images.append(img_path)

            slide = prs.slides.add_slide(prs.slide_layouts[6])
            page_ratio = meta['width'] / meta['height']
            slide_ratio = prs.slide_width / prs.slide_height
            if page_ratio >= slide_ratio:
                pic_w = prs.slide_width
                pic_h = int(pic_w / page_ratio)
                pic_left = 0
                pic_top = int((prs.slide_height - pic_h) / 2)
            else:
                pic_h = prs.slide_height
                pic_w = int(pic_h * page_ratio)
                pic_top = 0
                pic_left = int((prs.slide_width - pic_w) / 2)

            slide.shapes.add_picture(img_path, pic_left, pic_top, width=pic_w, height=pic_h)
            sx = pic_w / meta['width']
            sy = pic_h / meta['height']

            for line in meta['lines']:
                if len(line['text']) > 900 or not _region_is_near_white(image, line, meta['width'], meta['height']):
                    continue

                left = int(pic_left + line['x0'] * sx)
                top = int(pic_top + line['top'] * sy)
                width = max(int((line['x1'] - line['x0']) * sx) + 6, int(Inches(0.2)))
                height = max(int((line['bottom'] - line['top']) * sy * 1.35), int(Pt(line['size'] * 1.35)))

                cover = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left - 2, top - 1, width + 4, height + 2)
                cover.fill.solid()
                cover.fill.fore_color.rgb = RGBColor(255, 255, 255)
                cover.line.fill.background()

                textbox = slide.shapes.add_textbox(left, top, width, height)
                tf = textbox.text_frame
                tf.clear()
                tf.margin_left = 0
                tf.margin_right = 0
                tf.margin_top = 0
                tf.margin_bottom = 0
                tf.word_wrap = False
                tf.auto_size = MSO_AUTO_SIZE.NONE

                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = line['text']
                run.font.name = 'Arial'
                run.font.size = Pt(max(6.0, min(36.0, line['size'] * 0.98)))
                run.font.color.rgb = RGBColor(0, 0, 0)
                editable_line_count += 1

        output_file = os.path.join(output_dir, 'converted.pptx')
        prs.save(output_file)
        if not os.path.exists(output_file) or os.path.getsize(output_file) == 0:
            return {'success': False, 'error': 'PowerPoint file was not created'}
        if len(prs.slides) != page_count:
            return {
                'success': False,
                'error': f'PowerPoint validation failed: expected {page_count} slide(s), created {len(prs.slides)}.'
            }

        return {
            'success': True,
            'download_url': f'/download/{os.path.basename(output_dir)}/converted.pptx',
            'filename': 'converted.pptx',
            'message': (
                f'Converted {page_count} PDF page(s) to {page_count} PowerPoint slide(s). '
                f'Preserved the original page appearance and reconstructed {editable_line_count} text line(s) as editable content.'
            ),
        }
    except Exception as exc:
        return {'success': False, 'error': f'PDF to PowerPoint failed: {exc}'}
    finally:
        for img_path in temp_images:
            try:
                os.remove(img_path)
            except OSError:
                pass
